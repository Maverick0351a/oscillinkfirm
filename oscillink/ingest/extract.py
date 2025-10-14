from __future__ import annotations
# ruff: noqa: I001

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

import csv
import mimetypes
import mailbox
import hashlib
import re
from html.parser import HTMLParser
from email import policy
from email.parser import BytesParser
from email.utils import parsedate_to_datetime

import urllib.error
import urllib.parse
import urllib.request


@dataclass(frozen=True)
class ExtractedPage:
    page_number: int
    text: str
    source_path: str
    meta: Optional[Dict[str, Any]] = None


@dataclass(frozen=True)
class ExtractResult:
    pages: List[ExtractedPage]
    parser: str  # e.g., "tika", "pdfminer", "plain"
    source_path: str


def _extract_with_tika(path: Path, tika_url: str) -> Optional[List[ExtractedPage]]:
    """Call a Tika server to extract plain text. Returns pages or None on failure.

    We issue a PUT to the Tika endpoint with Accept: text/plain and treat the
    entire response as a single page. Uses stdlib urllib to avoid hard deps.
    """
    try:
        url = tika_url.rstrip("/")
        # Common Tika endpoints: "/tika" (auto-detect). We'll accept either a base or a full path.
        if not url.endswith("/tika"):
            url = url + "/tika"
        req = urllib.request.Request(
            url,
            data=path.read_bytes(),
            headers={"Accept": "text/plain"},
            method="PUT",
        )
        with urllib.request.urlopen(req, timeout=60) as resp:  # nosec - user-provided URL, documented feature
            data = resp.read()
        text = data.decode("utf-8", errors="replace")
        base = _file_meta(path, source_type=_source_type_from_suffix(path), mimetype=_guess_mimetype(path))
        meta = {**base, "page_or_row": "section:1"}
        pages = [ExtractedPage(page_number=1, text=text, source_path=str(path), meta=meta)]
        return pages
    except Exception:
        return None


def _extract_with_pdfminer(path: Path) -> Optional[List[ExtractedPage]]:
    """Use pdfminer.six when available to extract text. Returns pages or None on failure.

    We attempt to split pages using form-feed characters if present.
    """
    try:
        from pdfminer.high_level import extract_text  # type: ignore[import-not-found]

        text = extract_text(str(path))
        if text is None:
            return None
        # pdfminer often separates pages with \f
        parts = text.split("\f")
        pages: List[ExtractedPage] = []
        base = _file_meta(path, source_type="pdf", mimetype=_guess_mimetype(path) or "application/pdf")
        for i, t in enumerate(parts, start=1):
            t2 = t.strip()
            if not t2:
                continue
            meta = {**base, "page": i, "page_or_row": f"page:{i}"}
            pages.append(ExtractedPage(page_number=i, text=t2, source_path=str(path), meta=meta))
        if not pages and text:
            meta = {**base, "page": 1, "page_or_row": "page:1"}
            pages = [ExtractedPage(page_number=1, text=text, source_path=str(path), meta=meta)]
        return pages
    except Exception:
        return None


class _SimpleHTMLText(HTMLParser):
    """Minimal, deterministic HTML->text extraction preserving H1–H3 and paragraphs.

    - H1–H3 become markdown-style headings with '#'
    - <p> and <li> become paragraphs/lines
    - <pre>/<code> text is preserved
    - <title> is captured to self.title
    """

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._lines: List[str] = []
        self._cur_heading_level: Optional[int] = None
        self._in_p = False
        self._in_li = False
        self._in_pre = False
        self._in_title = False
        self.title: Optional[str] = None
        # shape hints
        self._in_table = False
        self._cur_row_cols = 0
        self.table_rows = 0
        self.table_cols = 0
        self.has_code = False

    def handle_starttag(self, tag: str, attrs):  # type: ignore[override]
        t = tag.lower()
        if t in ("h1", "h2", "h3"):
            self._cur_heading_level = int(t[1])
        elif t == "p":
            self._in_p = True
        elif t == "li":
            self._in_li = True
        elif t in ("pre", "code"):
            self._in_pre = True
            if t == "code":
                self.has_code = True
        elif t == "title":
            self._in_title = True
        elif t == "table":
            self._in_table = True
            self._cur_row_cols = 0
        elif t == "tr":
            self._cur_row_cols = 0
        elif t in ("td", "th") and self._in_table:
            self._cur_row_cols += 1

    def handle_endtag(self, tag: str):  # type: ignore[override]
        t = tag.lower()
        if t in ("h1", "h2", "h3"):
            # ensure separation after heading
            self._lines.append("")
            self._cur_heading_level = None
        elif t == "p":
            self._lines.append("")
            self._in_p = False
        elif t == "li":
            self._in_li = False
        elif t in ("pre", "code"):
            self._lines.append("")
            self._in_pre = False
        elif t == "title":
            self._in_title = False
        elif t == "tr" and self._in_table:
            self.table_rows += 1
            if self._cur_row_cols > self.table_cols:
                self.table_cols = self._cur_row_cols
        elif t == "table":
            self._in_table = False

    def handle_data(self, data: str):  # type: ignore[override]
        s = data.strip()
        if not s:
            return
        if self._in_title:
            self.title = (self.title or "") + (" " if self.title else "") + s
            return
        if self._cur_heading_level is not None:
            lvl = max(1, min(6, self._cur_heading_level))
            self._lines.append(f"{'#'*lvl} {s}")
            return
        if self._in_li:
            self._lines.append(f"- {s}")
            return
        if self._in_p or self._in_pre:
            self._lines.append(s)
            return
        # Fallback: inline text outside explicit blocks
        self._lines.append(s)

    def text(self) -> str:
        # Collapse multiple blank lines deterministically
        out: List[str] = []
        last_blank = False
        for line in self._lines:
            blank = line.strip() == ""
            if blank and last_blank:
                continue
            out.append(line)
            last_blank = blank
        return "\n".join(out).strip()


def _email_parts(msg) -> tuple[Optional[str], Optional[str], List[str]]:  # type: ignore[no-untyped-def]
    body_text: Optional[str] = None
    html_text: Optional[str] = None
    attachments: List[str] = []
    try:
        if msg.is_multipart():
            for part in msg.walk():
                cdisp = part.get_content_disposition()
                if cdisp == "attachment":
                    fn = part.get_filename()
                    if fn:
                        attachments.append(fn)
                    continue
                ctype = part.get_content_type()
                if body_text is None and ctype == "text/plain":
                    body_text = _get_part_text(part)
                elif html_text is None and ctype == "text/html":
                    html_text = _get_part_text(part)
        else:
            ctype = msg.get_content_type()
            if ctype == "text/plain":
                body_text = _get_part_text(msg)
            elif ctype == "text/html":
                html_text = _get_part_text(msg)
    except Exception:
        pass
    return body_text, html_text, attachments


def _get_part_text(part) -> Optional[str]:  # type: ignore[no-untyped-def]
    try:
        return part.get_content()
    except Exception:
        try:
            payload = part.get_payload(decode=True)
            if payload is None:
                return None
            charset = part.get_content_charset() or "utf-8"
            return payload.decode(charset, errors="replace")
        except Exception:
            return None


def _headers_from_msg(msg) -> Dict[str, Any]:  # type: ignore[no-untyped-def]
    meta: Dict[str, Any] = {}
    for k, key in (("subject", "subject"), ("from", "from"), ("to", "to"), ("date", "date"), ("message_id", "message-id")):
        try:
            val = msg.get(key)
        except Exception:
            val = None
        if val:
            meta[k] = str(val)
    return meta


def _header_prelude(meta: Dict[str, Any]) -> str:
    hdr_lines: List[str] = []
    subj = meta.get("subject")
    frm = meta.get("from")
    to = meta.get("to")
    if subj:
        hdr_lines.append(f"Subject: {subj}")
    if frm:
        hdr_lines.append(f"From: {frm}")
    if to:
        hdr_lines.append(f"To: {to}")
    return "\n".join(hdr_lines)


def _extract_html(path: Path) -> Optional[List[ExtractedPage]]:
    try:
        if not path.exists():
            return None
        text = path.read_text(encoding="utf-8", errors="replace")
        parser = _SimpleHTMLText()
        parser.feed(text)
        out_text = parser.text()
        base = _file_meta(path, source_type="html", mimetype=_guess_mimetype(path) or "text/html")
        meta = {**base, "page_or_row": "section:1"}
        if parser.title:
            meta["title"] = parser.title
        if parser.table_rows and parser.table_cols:
            meta["table_shape"] = {"rows": parser.table_rows, "cols": parser.table_cols}
        if parser.has_code:
            meta["has_code"] = True
        page = ExtractedPage(page_number=1, text=out_text, source_path=str(path), meta=meta)
        return [page]
    except Exception:
        return None


def _extract_eml(path: Path) -> Optional[List[ExtractedPage]]:
    """Extract text/metadata from an .eml message; return a single page.

    - Prefer text/plain body; fallback to text/html converted via _SimpleHTMLText
    - Attach headers and attachment filenames in metadata
    """
    try:
        if not path.exists():
            return None
        data = path.read_bytes()
        msg = BytesParser(policy=policy.default).parsebytes(data)
        body_text, html_text, attachments = _email_parts(msg)
        if body_text is None and html_text is not None:
            p = _SimpleHTMLText()
            p.feed(html_text)
            body_text = p.text()
        if body_text is None:
            body_text = ""
        base = _file_meta(path, source_type="eml", mimetype=_guess_mimetype(path) or "message/rfc822")
        meta: Dict[str, Any] = {**base, "page_or_row": "msg:1"}
        meta.update(_headers_from_msg(msg))
        # threading and seq
        _root, tid = _compute_thread_id(msg)
        if tid:
            meta["email.thread_id"] = tid
        if meta.get("message_id"):
            meta["email.message_id"] = meta["message_id"]
        meta["email.seq"] = 1
        if attachments:
            meta["attachments"] = attachments
        text_out = _header_prelude(meta) + ("\n\n" if meta.get("subject") else "") + body_text
        page = ExtractedPage(page_number=1, text=text_out, source_path=str(path), meta=meta)
        return [page]
    except Exception:
        return None


def _thread_sequence_for_mbox(parsed: List[tuple[int, Any]]) -> Dict[int, int]:
    """Compute per-thread sequence numbers for messages based on date, then index.

    Returns a mapping from message index (1-based position in mbox) to sequence number
    within its thread.
    """
    thread_map: Dict[str, List[tuple[int, Any, Optional[int]]]] = {}
    for idx, em in parsed:
        _root, tid = _compute_thread_id(em)
        key = tid or f"_no_thread_{idx}"
        ts = _parse_date(em)
        thread_map.setdefault(key, []).append((idx, em, ts))
    thread_seq: Dict[int, int] = {}
    for items in thread_map.values():
        items.sort(key=lambda t: (t[2] is None, t[2], t[0]))
        for i, (idx, _em, _ts) in enumerate(items, start=1):
            thread_seq[idx] = i
    return thread_seq


def _build_email_page_from_mbox(idx: int, em: Any, base: Dict[str, Any], src_path: Path, seq_map: Dict[int, int]) -> ExtractedPage:
    body_text, html_text, attachments = _email_parts(em)
    if body_text is None and html_text is not None:
        p = _SimpleHTMLText()
        p.feed(html_text)
        body_text = p.text()
    if body_text is None:
        body_text = ""
    meta: Dict[str, Any] = {**base, "message_index": idx, "page_or_row": f"msg:{idx}"}
    meta.update(_headers_from_msg(em))
    _root, tid = _compute_thread_id(em)
    if tid:
        meta["email.thread_id"] = tid
    if meta.get("message_id"):
        meta["email.message_id"] = meta["message_id"]
    if idx in seq_map:
        meta["email.seq"] = seq_map[idx]
    if attachments:
        meta["attachments"] = attachments
    text_out = _header_prelude(meta) + ("\n\n" if meta.get("subject") else "") + body_text
    return ExtractedPage(page_number=idx, text=text_out, source_path=str(src_path), meta=meta)


def _extract_mbox(path: Path) -> Optional[List[ExtractedPage]]:
    """Extract per-message pages from a .mbox mailbox file using stdlib mailbox."""
    try:
        if not path.exists():
            return None
        mbox = mailbox.mbox(str(path))
        base = _file_meta(path, source_type="mbox", mimetype=_guess_mimetype(path) or "application/mbox")
        # First pass: parse messages to email objects
        parsed: List[tuple[int, Any]] = []
        for idx, msg in enumerate(mbox, start=1):
            try:
                em = BytesParser(policy=policy.default).parsebytes(bytes(msg))
            except Exception:
                em = None  # type: ignore[assignment]
            if em is not None:
                parsed.append((idx, em))
        # Compute thread sequence mapping and build pages
        thread_seq = _thread_sequence_for_mbox(parsed)
        pages: List[ExtractedPage] = [
            _build_email_page_from_mbox(idx, em, base, path, thread_seq) for idx, em in parsed
        ]
        return pages or None
    except Exception:
        return None

def _extract_docx(path: Path) -> Optional[List[ExtractedPage]]:
    """Extract text from a DOCX file deterministically using python-docx if available.

    We preserve heading hierarchy by prefixing heading levels with markdown-style hashes.
    The entire document is returned as a single page; chunkers will split by paragraphs.
    """
    try:
        from docx import Document  # type: ignore[import-not-found]

        if not path.exists():
            return None
        doc = Document(str(path))
        lines: List[str] = []
        for para in doc.paragraphs:
            text = (para.text or "").strip()
            if not text:
                continue
            style = getattr(para, "style", None)
            name = getattr(style, "name", "") if style is not None else ""
            lvl = 0
            # Common heading names in python-docx: "Heading 1", "Heading 2", etc.
            if name and name.lower().startswith("heading"):
                try:
                    lvl = int(name.split()[-1])
                except Exception:
                    lvl = 1
            if lvl > 0:
                prefix = "#" * max(1, min(lvl, 6))
                lines.append(f"{prefix} {text}")
            else:
                lines.append(text)
        if not lines:
            return None
        joined = "\n\n".join(lines)
        base = _file_meta(path, source_type="docx", mimetype=_guess_mimetype(path) or "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        meta = {**base, "page_or_row": "section:1"}
        return [ExtractedPage(page_number=1, text=joined, source_path=str(path), meta=meta)]
    except Exception:
        return None


def _extract_pptx(path: Path) -> Optional[List[ExtractedPage]]:
    """Extract text per slide (and notes) from a PPTX using python-pptx if available.

    Returns one page per slide with page_number = slide index (1-based).
    """
    try:
        from pptx import Presentation  # type: ignore[import-not-found]

        if not path.exists():
            return None
        prs = Presentation(str(path))
        pages: List[ExtractedPage] = []
        base = _file_meta(path, source_type="pptx", mimetype=_guess_mimetype(path) or "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        for i, slide in enumerate(prs.slides, start=1):
            parts, has_notes = _pptx_slide_text(slide)
            text = "\n".join(parts).strip() or f"Slide {i}"
            pages.append(
                ExtractedPage(
                    page_number=i,
                    text=text,
                    source_path=str(path),
                    meta={**base, "slide_index": i, "has_notes": has_notes, "page_or_row": f"slide:{i}"},
                )
            )
        return pages if pages else None
    except Exception:
        return None


def _pptx_slide_text(slide) -> tuple[List[str], bool]:  # type: ignore[no-untyped-def]
    acc: List[str] = []
    has_notes = False
    for shape in getattr(slide, "shapes", []):
        try:
            if getattr(shape, "has_text_frame", False) and getattr(shape, "text_frame", None) is not None:
                for p in shape.text_frame.paragraphs:
                    t = (getattr(p, "text", "") or "").strip()
                    if t:
                        acc.append(t)
        except Exception:
            continue
    # Notes (optional)
    try:
        notes = getattr(slide, "notes_slide", None)
        ntf = getattr(notes, "notes_text_frame", None) if notes is not None else None
        if ntf is not None:
            for p in ntf.paragraphs:
                t = (getattr(p, "text", "") or "").strip()
                if t:
                    acc.append(f"[Notes] {t}")
                    has_notes = True
    except Exception:
        pass
    return acc, has_notes


def _render_row_text(headers: List[str], row: List[str], *, max_cols: int = 32, max_len: int = 2048) -> str:
    items: List[str] = []
    for h, v in zip(headers[:max_cols], row[:max_cols]):
        h2 = (h or "").strip()
        v2 = (str(v) if v is not None else "").strip()
        if h2 == "":
            items.append(v2)
        else:
            items.append(f"{h2}={v2}")
    s = "; ".join(items)
    if len(s) > max_len:
        s = s[: max_len - 3] + "..."
    return s


def _extract_csv(path: Path) -> Optional[List[ExtractedPage]]:
    """Extract deterministic row pages from CSV (first row as headers)."""
    try:
        with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
            rdr = csv.reader(f)
            headers: List[str] = []
            pages: List[ExtractedPage] = []
            base = _file_meta(path, source_type="csv", mimetype=_guess_mimetype(path) or "text/csv")
            for i, row in enumerate(rdr):
                if i == 0:
                    headers = [c.strip() for c in row]
                    continue
                text = _render_row_text(headers, [str(c) for c in row])
                # page_number = data row number (1-based)
                meta = {**base, "row_index": i, "page_or_row": f"row:{i}"}
                pages.append(ExtractedPage(page_number=i, text=text, source_path=str(path), meta=meta))
            return pages
    except Exception:
        return None


def _extract_xlsx(path: Path) -> Optional[List[ExtractedPage]]:
    """Extract deterministic row pages from XLSX using openpyxl if available.

    Produces one page per row per sheet. Page numbers are local to each sheet and
    continue across sheets to keep a total ordering.
    """
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]

        wb = load_workbook(filename=str(path), read_only=True, data_only=True)
        pages: List[ExtractedPage] = []
        page_counter = 0
        base = _file_meta(path, source_type="xlsx", mimetype=_guess_mimetype(path) or "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        for ws in wb.worksheets:
            rows_iter = ws.iter_rows(values_only=True)
            try:
                headers_row = next(rows_iter)
            except StopIteration:
                continue
            headers = [str(h) if h is not None else "" for h in headers_row]
            for sheet_row_index, row in enumerate(rows_iter, start=1):
                vals = [str(v) if v is not None else "" for v in row]
                txt = f"Sheet: {ws.title}\n" + _render_row_text(headers, vals)
                page_counter += 1
                pages.append(
                    ExtractedPage(
                        page_number=page_counter,
                        text=txt,
                        source_path=str(path),
                        meta={**base, "sheet": ws.title, "sheet_row_index": sheet_row_index, "page_or_row": f"sheet:{ws.title}!row:{sheet_row_index}"},
                    )
                )
        return pages if pages else None
    except Exception:
        return None


def extract_text(paths: Iterable[str | Path], *, parser: str = "plain", tika_url: Optional[str] = None) -> List[ExtractResult]:
    """Deterministic text extraction stub.

    For now, provides a trivial, deterministic implementation:
    - If a path is a .txt file, read as a single page.
    - If parser is 'tika' and a tika_url is provided, call Tika server.
    - If parser is 'pdfminer' and pdfminer.six is installed, use it.
    - If parser is 'auto', try txt -> tika -> pdfminer, then fallback empty.
    - Otherwise, return empty result; OCR step may handle it.

    This keeps the repo green until heavy deps (Tika, pdfminer) are added.
    """
    results: List[ExtractResult] = []
    for p in paths:
        sp = str(p)
        path = Path(p)
        used_parser = parser
        pages: List[ExtractedPage] | None = None

        # Always handle simple txt/md deterministically
        if path.suffix.lower() in (".txt", ".md") and path.exists():
            txt = path.read_text(encoding="utf-8", errors="replace")
            base = _file_meta(path, source_type=_source_type_from_suffix(path), mimetype=_guess_mimetype(path) or ("text/markdown" if path.suffix.lower() == ".md" else "text/plain"))
            meta = {**base, "page_or_row": "section:1"}
            pages = [ExtractedPage(page_number=1, text=txt, source_path=sp, meta=meta)]
            used_parser = "plain"
        else:
            pages, tag = _extract_auto_for_path(path, tika_url) if parser == "auto" else _extract_with_named_parser(path, parser, tika_url)
            if pages is not None and tag:
                used_parser = tag

        if pages is None:
            pages = []
        results.append(ExtractResult(pages=pages, parser=used_parser, source_path=sp))
    return results


def _guess_mimetype(path: Path) -> Optional[str]:
    try:
        mtype, _ = mimetypes.guess_type(str(path))
        return mtype
    except Exception:
        return None


def _source_type_from_suffix(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    return ext or "file"


def _file_meta(path: Path, *, source_type: str, mimetype: Optional[str]) -> Dict[str, Any]:
    """Build uniform file-level metadata.

    - created_at, modified_at: integer epoch seconds
    - mimetype: best-effort guess
    - source_type: normalized from extension (e.g., pdf, docx, xlsx)
    """
    try:
        st = path.stat()
        created = int(getattr(st, "st_ctime", getattr(st, "st_mtime", 0)))
        modified = int(getattr(st, "st_mtime", 0))
    except Exception:
        created = 0
        modified = 0
    out: Dict[str, Any] = {
        "created_at": created,
        "modified_at": modified,
        "source_type": source_type,
    }
    if mimetype:
        out["mimetype"] = mimetype
    # Optional education folder mapping: courses/{COURSE}/{SEM}/Week{N}/...
    try:
        parts = list(path.parts)
        for i, seg in enumerate(parts):
            if seg.lower() == "courses" and i + 2 < len(parts):
                course = parts[i + 1]
                sem = parts[i + 2]
                out["edu.course_code"] = course
                out["edu.semester"] = sem
                # find a WeekN segment
                for s in parts[i + 3 : i + 8]:  # look ahead a few segments
                    m = re.match(r"(?i)week\s*([0-9]+)", s)
                    if m:
                        out["edu.week"] = int(m.group(1))
                        break
                break
    except Exception:
        pass
    return out


def _compute_thread_id(msg) -> tuple[Optional[str], Optional[str]]:  # type: ignore[no-untyped-def]
    try:
        in_reply = msg.get("In-Reply-To") or msg.get("in-reply-to")
        refs = msg.get("References") or msg.get("references")
        msgid = msg.get("Message-ID") or msg.get("message-id")
        root = None
        if in_reply:
            root = str(in_reply)
        elif refs:
            # take first reference token
            root = str(refs).split()[0]
        elif msgid:
            root = str(msgid)
        if not root:
            return None, None
        h = hashlib.sha256(root.encode("utf-8", errors="ignore")).hexdigest()[:16]
        return root, h
    except Exception:
        return None, None


def _parse_date(msg) -> Optional[int]:  # type: ignore[no-untyped-def]
    try:
        d = msg.get("Date") or msg.get("date")
        if not d:
            return None
        dt = parsedate_to_datetime(str(d))
        if dt is None:
            return None
        return int(dt.timestamp())
    except Exception:
        return None


def _extract_auto_for_path(path: Path, tika_url: Optional[str]) -> tuple[Optional[List[ExtractedPage]], Optional[str]]:
    pages: Optional[List[ExtractedPage]] = None
    tag: Optional[str] = None
    ext = path.suffix.lower()
    # Try direct by extension first
    ext_map = {
        ".docx": ("docx", _extract_docx),
        ".pptx": ("pptx", _extract_pptx),
        ".csv": ("csv", _extract_csv),
        ".xlsx": ("xlsx", _extract_xlsx),
        ".html": ("html", _extract_html),
        ".htm": ("html", _extract_html),
        ".eml": ("eml", _extract_eml),
        ".mbox": ("mbox", _extract_mbox),
    }
    if ext in ext_map:
        tag0, fn = ext_map[ext]
        pages = fn(path)
        if pages is not None:
            tag = tag0
            return pages, tag
    # Try tika when available
    if tika_url:
        pages = _extract_with_tika(path, tika_url)
        if pages is not None:
            return pages, "tika"
    # Try pdfminer for PDFs
    if ext == ".pdf":
        pages = _extract_with_pdfminer(path)
        if pages is not None:
            return pages, "pdfminer"
    return None, None


def _extract_with_named_parser(path: Path, parser: str, tika_url: Optional[str]) -> tuple[Optional[List[ExtractedPage]], Optional[str]]:
    if parser == "tika":
        return (_extract_with_tika(path, tika_url) if tika_url else None), ("tika" if tika_url else None)
    parser_map = {
        "pdfminer": _extract_with_pdfminer,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
        "csv": _extract_csv,
        "xlsx": _extract_xlsx,
        "html": _extract_html,
        "eml": _extract_eml,
        "mbox": _extract_mbox,
    }
    fn = parser_map.get(parser)
    if fn is None:
        return None, None
    pages = fn(path)
    return pages, (parser if pages is not None else None)
